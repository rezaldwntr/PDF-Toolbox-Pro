# app/routers/convert.py
import os
import shutil
import logging
import tempfile
from zipfile import ZipFile

from fastapi import APIRouter, File, UploadFile, HTTPException, BackgroundTasks
from fastapi.responses import FileResponse

# Library Konversi
from pdf2docx import Converter
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
import pdfplumber
import pandas as pd
from openpyxl.styles import Border, Side, Alignment, Font, PatternFill

from app.utils.file_utils import validate_file, cleanup_folder

router = APIRouter(prefix="/convert", tags=["Conversion"])

# --- HELPER KHUSUS EXCEL ---
def dataframe_to_rows(df, index=False, header=False):
    if header: yield df.columns.tolist()
    for row in df.itertuples(index=index, name=None): yield row

def is_inside(bbox, tables):
    mx, my = (bbox[0]+bbox[2])/2, (bbox[1]+bbox[3])/2
    for t in tables:
        if t[0]<=mx<=t[2] and t[1]<=my<=t[3]: return True
    return False

# === 1. PDF KE DOCX ===
@router.post("/pdf-to-docx")
def convert_pdf_to_docx(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    validate_file(file)
    tmp_dir = tempfile.mkdtemp()
    tmp_pdf_path = os.path.join(tmp_dir, file.filename)
    docx_filename = os.path.splitext(file.filename)[0] + ".docx"
    tmp_docx_path = os.path.join(tmp_dir, docx_filename)
    try:
        with open(tmp_pdf_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
        
        cv = Converter(tmp_pdf_path)
        cv.convert(tmp_docx_path, start=0, end=None, multiprocess=False)
        cv.close()
        
        background_tasks.add_task(cleanup_folder, tmp_dir)
        return FileResponse(path=tmp_docx_path, filename=docx_filename, media_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document')
    except Exception as e:
        cleanup_folder(tmp_dir)
        raise HTTPException(status_code=500, detail=f"Gagal convert Word: {str(e)}")

# === 2. PDF KE EXCEL ===
@router.post("/pdf-to-excel")
def convert_pdf_to_excel(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    validate_file(file)
    tmp_dir = tempfile.mkdtemp()
    tmp_pdf_path = os.path.join(tmp_dir, file.filename)
    xlsx_filename = os.path.splitext(file.filename)[0] + ".xlsx"
    tmp_xlsx_path = os.path.join(tmp_dir, xlsx_filename)
    try:
        with open(tmp_pdf_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
        
        with pd.ExcelWriter(tmp_xlsx_path, engine='openpyxl') as writer:
            with pdfplumber.open(tmp_pdf_path) as pdf:
                writer.book.create_sheet("Hasil Konversi")
                worksheet = writer.book["Hasil Konversi"]
                if "Sheet" in writer.book.sheetnames: del writer.book["Sheet"]
                
                current_row = 1
                thin_border = Border(left=Side(style='thin'), right=Side(style='thin'), top=Side(style='thin'), bottom=Side(style='thin'))
                
                for page in pdf.pages:
                    tables = page.find_tables()
                    table_bboxes = [t.bbox for t in tables]
                    words = page.extract_words()
                    non_table_words = [w for w in words if not is_inside((w['x0'],w['top'],w['x1'],w['bottom']), table_bboxes)]
                    
                    non_table_text = []
                    if non_table_words:
                         non_table_text = sorted(non_table_words, key=lambda x: x['top'])
                    
                    for w in non_table_text:
                        worksheet.cell(row=current_row, column=1, value=w['text'])
                        current_row += 1
                    
                    for t in tables:
                        data = t.extract()
                        if data:
                            df = pd.DataFrame(data)
                            for r in dataframe_to_rows(df, index=False, header=False):
                                for c_idx, val in enumerate(r, 1):
                                    c = worksheet.cell(row=current_row, column=c_idx, value=val)
                                    c.border = thin_border
                                current_row += 1
                        current_row += 1
                    current_row += 2
        
        background_tasks.add_task(cleanup_folder, tmp_dir)
        return FileResponse(path=tmp_xlsx_path, filename=xlsx_filename, media_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
    except Exception as e:
        cleanup_folder(tmp_dir)
        logging.error(f"Excel error: {e}")
        raise HTTPException(status_code=500, detail="Gagal convert Excel. Pastikan file tidak corrupt.")

# === 3. PDF KE PPT ===
@router.post("/pdf-to-ppt")
def convert_pdf_to_ppt(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    validate_file(file)
    tmp_dir = tempfile.mkdtemp()
    tmp_pdf_path = os.path.join(tmp_dir, file.filename)
    ppt_filename = os.path.splitext(file.filename)[0] + ".pptx"
    tmp_ppt_path = os.path.join(tmp_dir, ppt_filename)
    try:
        with open(tmp_pdf_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
        
        prs = Presentation()
        doc = fitz.open(tmp_pdf_path)
        if len(doc) > 0:
            p1 = doc[0]
            prs.slide_width = int((p1.rect.width / 72) * 914400)
            prs.slide_height = int((p1.rect.height / 72) * 914400)
        
        for page in doc:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            # Images
            img_blocks = [b for b in page.get_text("dict", flags=fitz.TEXT_PRESERVE_IMAGES)["blocks"] if b['type']==1]
            for b in img_blocks:
                img_path = os.path.join(tmp_dir, f"img_{os.urandom(4).hex()}.{b['ext']}")
                with open(img_path, "wb") as f: f.write(b["image"])
                x0,y0,x1,y1 = b["bbox"]
                try: slide.shapes.add_picture(img_path, Inches(x0/72), Inches(y0/72), width=Inches((x1-x0)/72), height=Inches((y1-y0)/72))
                except: pass
            
            # Text
            text_blocks = [b for b in page.get_text("dict")["blocks"] if b['type']==0]
            for b in text_blocks:
                for line in b["lines"]:
                    lx0,ly0,lx1,ly1 = line["bbox"]
                    txBox = slide.shapes.add_textbox(Inches(lx0/72), Inches(ly0/72), Inches((lx1-lx0)/72), Inches((ly1-ly0)/72))
                    tf = txBox.text_frame
                    tf.word_wrap = False
                    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
                    p = tf.paragraphs[0]
                    for span in line["spans"]:
                        if not span["text"].strip(): continue
                        run = p.add_run()
                        run.text = span["text"]
                        run.font.size = Pt(span["size"])
                        try:
                            c = span["color"]
                            run.font.color.rgb = RGBColor((c>>16)&0xFF, (c>>8)&0xFF, c&0xFF)
                        except: pass
                        if span["flags"] & 16: run.font.bold = True
                        if span["flags"] & 2: run.font.italic = True
        doc.close()
        prs.save(tmp_ppt_path)
        background_tasks.add_task(cleanup_folder, tmp_dir)
        return FileResponse(path=tmp_ppt_path, filename=ppt_filename, media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    except Exception as e:
        cleanup_folder(tmp_dir)
        raise HTTPException(status_code=500, detail=f"Gagal convert PPT: {str(e)}")

# === 4. PDF KE IMAGE ===
@router.post("/pdf-to-image")
def convert_pdf_to_image(background_tasks: BackgroundTasks, file: UploadFile = File(...), output_format: str = "png"):
    validate_file(file)
    tmp_dir = tempfile.mkdtemp()
    tmp_pdf_path = os.path.join(tmp_dir, file.filename)
    zip_filename = os.path.splitext(file.filename)[0] + "_images.zip"
    tmp_zip_path = os.path.join(tmp_dir, zip_filename)
    try:
        with open(tmp_pdf_path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
        
        doc = fitz.open(tmp_pdf_path)
        with ZipFile(tmp_zip_path, 'w') as zipf:
            for i, page in enumerate(doc):
                pix = page.get_pixmap(dpi=200)
                img_name = f"page_{i+1}.{output_format}"
                img_path = os.path.join(tmp_dir, img_name)
                pix.save(img_path, output="jpg" if output_format.lower() in ['jpg','jpeg'] else output_format)
                zipf.write(img_path, img_name)
        doc.close()
        background_tasks.add_task(cleanup_folder, tmp_dir)
        return FileResponse(path=tmp_zip_path, filename=zip_filename, media_type='application/zip')
    except Exception as e:
        cleanup_folder(tmp_dir)
        raise HTTPException(status_code=500, detail=f"Gagal convert Image: {str(e)}")
